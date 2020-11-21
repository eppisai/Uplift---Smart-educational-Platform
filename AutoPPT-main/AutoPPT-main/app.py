# Importing flask module in the project is mandatory
# An object of Flask class is our WSGI application.
from flask import Flask, redirect
from flask import request

# Web scraping, pptx, requests imports                                                                                     @****
import requests
from bs4 import BeautifulSoup
from pptx.util import Inches, Pt
from pptx import Presentation
import urllib.request


# Flask constructor takes the name of
# current module (__name__) as argument.
app = Flask(__name__)

# The route() function of the Flask class is a decorator,
# which tells the application which URL should call
# the associated function.


@app.route('/flask/')
# ‘/’ URL is bound with hello_world() function.
def hello_world():
    args = request.args
    # topic = input("what is the topic !? :)")  # @***
    topic = args['topic']

    # user = input("what's your name?")
    user = args['user']

    # this line uses 'get' function to make a response object and to get the source code , use .text @***
    try:
        source = requests.get('https://en.wikipedia.org/wiki/' + topic).text
    except:
        print("Error opening the URL")

        
        

    soup = BeautifulSoup(source, 'lxml')
    # the next line gives the entire source code                                                                          @****
    # print(soup.prettify()) <--------------this is source code!

    # one seperate soup of headings as well as paragraphs, this helps in targeting the topics individually, later.
    article = " "
    for heading in soup.find_all(["p", "h2", "h3"]):
        article += heading.name + ' ' + heading.text.strip()

    def getdata(url):
     r = requests.get(url)
     return r.text


    htmldata = getdata("https://www.geeksforgeeks.org/")
    soup2 = BeautifulSoup(htmldata, 'html.parser')
    for item in soup2.find_all('img'):
         try:
            source = requests.get(item['src'])
         except:
            print("Error opening the URL")

         file = open("sample_image.png", "wb")
         file.write(source.content)
         file.close()

    # ..                                                                                                                     ****
    first_slide = ""
    second_slide = ""
    third_slide = ""
    fourth_slide = ""
    fifth_slide = ""
    sixth_slide = ""
    seventh_slide = ""
    matter = ''
    deffi_slides = []
    count = 0
    lines = 0
    '''
    this piece of code is specific to wikipidea.
    every page has its last hyperlinked heading as "see also" and we dont want it.
    also, the index was too big to fit into one slide, so i split the entries in more slides.the "if" section is for that slpitting.
    below it, we have 4 lines of code which extract out the first line of each heading by finding  the first full stop available.
    '''
    for heading in soup.find_all(["h2", "h3"]):
        if heading.text.strip() == 'See also':
            break
        if heading.name == "h2":
            char = ' ' + heading.text.strip()+"\n"
        if heading.name == "h3":
            char = '  -- ' + heading.text.strip() + "\n"
        lines += 1
        matter = ''
        # -----------------divides the *contents* section into parts of 9 lines per slide. to avoid congession

        if lines <= 9 and (not char.isdigit()):
            first_slide = first_slide+char
            count = 1
        if lines > 9 and lines < +18:
            second_slide = second_slide+char
            count = 2
        if lines > 18 and lines <= 27:
            third_slide = third_slide+char
            count = 3
        if lines > 27 and lines <= 36:
            fourth_slide = fourth_slide+char
            count = 4
        if lines > 36 and lines <= 45:
            fifth_slide = fifth_slide+char
            count = 5
        if lines > 45 and lines <= 54:
            sixth_slide = sixth_slide+char
            count = 6
        if lines > 54 and lines <= 63:
            seventh_slide = seventh_slide+char
            count = 7
            # ----------------------------------- finding the first sentences---------------
        start_def = article.find(heading.text.strip())
        end_def = article.find('.', start_def)
        matter += article[start_def: end_def+1]
        deffi_slides = deffi_slides + \
            [[str(heading.text.strip()), str(matter), 1]]

    prs = Presentation()
    # the format of this class is to refer to a 'list if lists'(make sure you remember this,or  it can lead to a frustrating mistake).
    # i found the class format, as the only way to add a second slide.so every required data is put into that format

    class MySlide:
        def __init__(self, data):
            self.layout = prs.slide_layouts[data[2]]
            # apply the layout to the slide
            self.slide = prs.slides.add_slide(self.layout)
            self.title = self.slide.shapes.title
            self.title.text = data[0]
            self.subtitle = self.slide.placeholders[1]
            self.subtitle.text = data[1]

    slides = [[soup.h1.text, " by " + user+"", 0]]
    index_slide = [first_slide, second_slide, third_slide,
                   fourth_slide, fifth_slide, sixth_slide, seventh_slide]
    indexing = []
    for a in range(count):
        indexing = indexing+[["contents", index_slide[a], 1]]
    # slides= slides+indexing and so on.... slides should be the final, only list, for easier addition.

    slides = slides + indexing + deffi_slides

    for es in slides:
        MySlide(es)

    prs.save('something.pptx')

    return redirect("http://localhost:3000/flask", code=302)


# main driver function
if __name__ == '__main__':

    # run() method of Flask class runs the application
    # on the local development server.
    app.run()
